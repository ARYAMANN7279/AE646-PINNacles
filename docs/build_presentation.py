"""
PINNacles Stage 1 Proposal Presentation - v3
All positions hardcoded; no computed remainder heights.
Safe content zone: top=1.55, bottom=7.05 (never exceeded).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE   = RGBColor(0x1A, 0x2E, 0x4A)
MID_BLUE    = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF4, 0xF6, 0xF9)
DARK_GRAY   = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY    = RGBColor(0x5D, 0x6D, 0x7E)
GREEN       = RGBColor(0x1E, 0x8B, 0x4C)
ORANGE      = RGBColor(0xCA, 0x6F, 0x1E)

W, H = Inches(13.33), Inches(7.5)
# Safe content zone
CL, CR = Inches(0.45), Inches(12.88)   # left / right margin edge
CW = CR - CL                            # 12.43"
CT = Inches(1.55)                       # below header bar + gap
CB = Inches(7.05)                       # above footer + gap
CAH = CB - CT                           # 5.50" usable

TOTAL = 10

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── primitives ────────────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(BLANK)

def box(sl, l, t, w, h, fill=None, border=None, bw=Pt(1.2)):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.line.width = bw
    if fill:  s.fill.solid(); s.fill.fore_color.rgb = fill
    else:     s.fill.background()
    if border: s.line.color.rgb = border
    else:      s.line.fill.background()
    return s

def tx(sl, text, l, t, w, h, sz=Pt(16), bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def blist(sl, items, l, t, w, h, sz=Pt(16), color=DARK_GRAY, spacing=Pt(5)):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        lvl = 1 if item.startswith("  ") else 0
        p.level = lvl
        r = p.add_run()
        r.text = ("◦  " if lvl else "•  ") + item.strip()
        r.font.size = sz; r.font.color.rgb = color; r.font.name = "Calibri"

def header(sl, title, sub=None):
    box(sl, 0, 0, W, Inches(1.35), fill=DARK_BLUE)
    box(sl, 0, Inches(1.35), W, Inches(0.055), fill=ACCENT_BLUE)
    tx(sl, title, Inches(0.45), Inches(0.1), W-Inches(0.9), Inches(0.75),
       sz=Pt(30), bold=True, color=WHITE)
    if sub:
        tx(sl, sub, Inches(0.45), Inches(0.86), W-Inches(0.9), Inches(0.42),
           sz=Pt(15), color=LIGHT_BLUE)

def footer(sl, n):
    box(sl, 0, Inches(7.12), W, Inches(0.38), fill=DARK_BLUE)
    tx(sl, "AE646  |  PINNacles  |  Stage 1 Proposal",
       Inches(0.4), Inches(7.16), Inches(9), Inches(0.3), sz=Pt(11), color=LIGHT_BLUE)
    tx(sl, f"{n} / {TOTAL}", W-Inches(1.1), Inches(7.16), Inches(0.9), Inches(0.3),
       sz=Pt(11), color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

def card_hdr(sl, label, l, t, w, h=Inches(0.48), fill=ACCENT_BLUE):
    box(sl, l, t, w, h, fill=fill)
    tx(sl, label, l+Inches(0.15), t+Inches(0.08), w-Inches(0.3), h-Inches(0.1),
       sz=Pt(17), bold=True, color=WHITE)

def card_body(sl, l, t, w, h):
    box(sl, l, t, w, h, fill=WHITE, border=LIGHT_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# S1 Title
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=DARK_BLUE)
box(sl, 0, Inches(3.55), W, Inches(0.07), fill=ACCENT_BLUE)
tx(sl, "Fourier Neural Operator\nfor Parametric Darcy Flow",
   Inches(1.0), Inches(0.55), W-Inches(2.0), Inches(2.75),
   sz=Pt(46), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "AE646: Scientific Machine Learning for Fluid Mechanics\nProject Theme #8 - Operator Learning for Parametric PDEs",
   Inches(1.0), Inches(3.75), W-Inches(2.0), Inches(0.9),
   sz=Pt(19), color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
tx(sl, "Team PINNacles",
   Inches(1.0), Inches(4.82), W-Inches(2.0), Inches(0.48),
   sz=Pt(21), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
tx(sl, "Aryamann Srivastava  •  Varun Sathaye  •  Atishay Jain  •  Vedant S. Tiwari",
   Inches(1.0), Inches(5.35), W-Inches(2.0), Inches(0.42),
   sz=Pt(17), color=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "IIT Kanpur  |  Semester 2026–27-I  |  Dr. Sathesh Mariappan",
   Inches(1.0), Inches(5.9), W-Inches(2.0), Inches(0.36),
   sz=Pt(13), color=MED_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# S2 Problem Statement
# layout: [PDE box left | Challenge box right] top row 2.1"
#         [Objective banner] bottom row 3.1"
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Problem Statement",
       "2D Darcy flow - steady-state pressure in a heterogeneous porous medium")
footer(sl, 2)

HALF = (CW - Inches(0.2)) / 2

# PDE box  t=1.55 h=2.1 b=3.65
card_body(sl, CL, CT, HALF, Inches(2.1))
tx(sl, "The Governing Equation", CL+Inches(0.18), CT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
tx(sl, "−∇·( κ(x,y) ∇u(x,y) ) = f(x,y)   on (0,1)²",
   CL+Inches(0.18), CT+Inches(0.57), HALF-Inches(0.36), Inches(0.48),
   sz=Pt(17), bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
tx(sl, "u = 0 on boundary,   f = β = 1.0",
   CL+Inches(0.18), CT+Inches(1.1), HALF-Inches(0.36), Inches(0.36),
   sz=Pt(15), color=DARK_GRAY, align=PP_ALIGN.CENTER)
tx(sl, "κ(x,y) = permeability field   |   u(x,y) = pressure solution",
   CL+Inches(0.18), CT+Inches(1.52), HALF-Inches(0.36), Inches(0.48),
   sz=Pt(14), color=MED_GRAY, italic=True, align=PP_ALIGN.CENTER)

# Challenge box
RX = CL + HALF + Inches(0.2)
card_body(sl, RX, CT, HALF, Inches(2.1))
tx(sl, "The Challenge", RX+Inches(0.18), CT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=ORANGE)
blist(sl,
    ["Traditional solvers (FDM/FEM) re-solve for every new κ field",
     "Each solve: ~1 second for a simple 2D case",
     "Optimization or UQ needs thousands of solves → impractical"],
    RX, CT+Inches(0.57), HALF, Inches(1.45), sz=Pt(16))

# Objective banner  t=3.75 h=3.3 b=7.05
OT = CT + Inches(2.1) + Inches(0.1)
box(sl, CL, OT, CW, Inches(3.2), fill=LIGHT_BLUE, border=ACCENT_BLUE, bw=Pt(2))
tx(sl, "Objective", CL+Inches(0.25), OT+Inches(0.18), Inches(2.5), Inches(0.45),
   sz=Pt(22), bold=True, color=MID_BLUE)
tx(sl, "Learn the solution operator  G : κ(x,y) → u(x,y)  once.\n\n"
       "Evaluate any new permeability field in milliseconds - no re-solving needed.\n\n"
       "Enables fast parametric studies, optimization, and uncertainty quantification.",
   CL+Inches(0.25), OT+Inches(0.72), CW-Inches(0.5), Inches(2.35),
   sz=Pt(18), color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# S3 Dataset
# layout: [Source | Splits] top row 2.1"
#         [κ field | Why 64×64 | Super-res] bottom row 2.85"
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Dataset", "PDEBench 2D Darcy Flow (β=1.0) - official, checksum-verified source")
footer(sl, 3)

card_body(sl, CL, CT, HALF, Inches(2.1))
tx(sl, "Official Source", CL+Inches(0.18), CT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
tx(sl, "DaRUS / University of Stuttgart\nDOI: 10.18419/darus-2986\n10,000 samples at native 128×128 resolution",
   CL+Inches(0.18), CT+Inches(0.58), HALF-Inches(0.36), Inches(1.42),
   sz=Pt(16), color=DARK_GRAY)

card_body(sl, RX, CT, HALF, Inches(2.1))
tx(sl, "Our Reproducible Subset  (seed = 42)", RX+Inches(0.18), CT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
tx(sl, "Train: 900 samples\nValidation: 100 samples\nTest: 200 samples (held out)\nMain resolution: 64×64  (stride-2 downsample)",
   RX+Inches(0.18), CT+Inches(0.58), HALF-Inches(0.36), Inches(1.42),
   sz=Pt(16), color=DARK_GRAY)

# 3 bottom cards  t=3.75 h=3.3 b=7.05
C3W = (CW - Inches(0.4)) / 3
BT3 = CT + Inches(2.1) + Inches(0.1)
BH3 = Inches(3.2)
for i, (label, body, col) in enumerate([
    ("Permeability  κ(x,y)",
     "Piecewise-constant: values in {0.1, 1.0}\n\nGenerated by thresholding a smooth Gaussian random field at zero\n\nProduces sharp-boundary blob patterns in 2D",
     ACCENT_BLUE),
    ("Why 64×64?",
     "Standard resolution for FNO Darcy experiments\n\nRetains full PDE spatial structure\n\nFNO can then be tested at 128×128 zero-shot - no retraining needed",
     GREEN),
    ("Super-Res Hold-out",
     "200 test samples also kept at native 128×128\n\nUsed only for zero-shot evaluation\n\nNever seen by the model during training in any form",
     ORANGE),
]):
    CX3 = CL + i * (C3W + Inches(0.2))
    card_hdr(sl, label, CX3, BT3, C3W, fill=col)
    card_body(sl, CX3, BT3+Inches(0.48), C3W, BH3-Inches(0.48))
    tx(sl, body, CX3+Inches(0.18), BT3+Inches(0.62), C3W-Inches(0.36), BH3-Inches(0.75),
       sz=Pt(16), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# S4 Baseline: MLP
# layout: [flow strip] 1.1"
#         [Architecture details | Why baseline] 2.75"
#         [note banner] 1.2"
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Baseline Method: MLP", "Dense fully-connected network - no spatial inductive bias")
footer(sl, 4)

# Flow strip
FSH = Inches(1.1)
card_body(sl, CL, CT, CW, FSH)
blks = [
    ("Input\n64×64×3", ACCENT_BLUE),
    ("Flatten\n→12,288", MID_BLUE),
    ("Linear 2048\n+GELU", MID_BLUE),
    ("Linear 2048\n+GELU", MID_BLUE),
    ("Linear 2048\n+GELU", MID_BLUE),
    ("Linear 4096\n+Reshape", MID_BLUE),
    ("Output\n64×64×1", GREEN),
]
AW = Inches(0.25)
BW = (CW - Inches(0.36) - AW * (len(blks)-1)) / len(blks)
BH = Inches(0.78); BT_ = CT + (FSH - BH)/2
for j, (lbl, col) in enumerate(blks):
    bx = CL + Inches(0.18) + j*(BW+AW)
    box(sl, bx, BT_, BW, BH, fill=col)
    tx(sl, lbl, bx, BT_+Inches(0.06), BW, BH-Inches(0.06),
       sz=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if j < len(blks)-1:
        tx(sl, "→", bx+BW, BT_+Inches(0.22), AW, Inches(0.3),
           sz=Pt(16), bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Two detail cards  t=2.75 h=2.75 b=5.5
DT = CT + FSH + Inches(0.1)
DH = Inches(2.75)
card_body(sl, CL, DT, HALF, DH)
tx(sl, "Architecture Details", CL+Inches(0.18), DT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
blist(sl,
    ["42.0M trainable parameters",
     "Input: [permeability, x_coord, y_coord] - 3 channels",
     "3 hidden layers × 2048 units, GELU activation",
     "Output: flat 64×64, reshaped to grid"],
    CL, DT+Inches(0.58), HALF, DH-Inches(0.65), sz=Pt(16))

card_body(sl, RX, DT, HALF, DH)
tx(sl, "Why This as Baseline?", RX+Inches(0.18), DT+Inches(0.12), HALF-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=ORANGE)
blist(sl,
    ["Treats all 4096 output pixels independently",
     "No spatial locality or translation equivariance",
     "Fixed 64×64 I/O - cannot evaluate at any other resolution",
     "Establishes a clear performance floor for comparison"],
    RX, DT+Inches(0.58), HALF, DH-Inches(0.65), sz=Pt(16))

# Note banner  t=5.6 h=1.45 b=7.05
NT = DT + DH + Inches(0.1)
box(sl, CL, NT, CW, Inches(1.45), fill=LIGHT_BLUE, border=ACCENT_BLUE, bw=Pt(2))
tx(sl, "Key limitation:  the MLP has hard-coded 64×64 input/output dimensions - it "
       "cannot be evaluated zero-shot at any other resolution. This is in direct contrast to the FNO.",
   CL+Inches(0.25), NT+Inches(0.22), CW-Inches(0.5), Inches(1.0),
   sz=Pt(17), color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# S5 Proposed Method: FNO
# layout: [flow strip] 1.1"
#         [How SpectralConv works | Key Properties] 4.05"
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Proposed SciML Method: Fourier Neural Operator",
       "Operator learning via spectral convolutions - resolution-invariant by construction")
footer(sl, 5)

card_body(sl, CL, CT, CW, FSH)
fblks = [
    ("Input\n64×64×3", ACCENT_BLUE),
    ("Lift\n3→64", MID_BLUE),
    ("SpectralConv\nBlock ×4", MID_BLUE),
    ("FFT→W→IFFT\n+Conv1×1\n+LN+GELU", MID_BLUE),
    ("Project\n64→128→1", MID_BLUE),
    ("Output\n64×64×1", GREEN),
]
AW5 = Inches(0.25)
FBW = (CW - Inches(0.36) - AW5 * (len(fblks)-1)) / len(fblks)
for j, (lbl, col) in enumerate(fblks):
    bx = CL + Inches(0.18) + j*(FBW+AW5)
    box(sl, bx, BT_, FBW, BH, fill=col)
    tx(sl, lbl, bx, BT_+Inches(0.04), FBW, BH-Inches(0.04),
       sz=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if j < len(fblks)-1:
        tx(sl, "→", bx+FBW, BT_+Inches(0.22), AW5, Inches(0.3),
           sz=Pt(16), bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Two bottom cards  t=2.75 h=4.3 b=7.05
BT5 = CT + FSH + Inches(0.1)
BH5 = Inches(4.3)
SW  = CW * 0.52
PW  = CW - SW - Inches(0.2)

card_body(sl, CL, BT5, SW, BH5)
tx(sl, "How Spectral Convolution Works",
   CL+Inches(0.18), BT5+Inches(0.12), SW-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
steps = [
    "Apply 2D FFT  →  frequency domain representation",
    "Keep only the lowest 12×12 Fourier modes",
    "Multiply by learned complex weights (learning in Fourier space)",
    "Inverse FFT  →  back to physical space",
    "Add 1×1 Conv bypass, normalize, activate (GELU)",
]
for k, s in enumerate(steps):
    tx(sl, f"{k+1}.  {s}",
       CL+Inches(0.25), BT5+Inches(0.65)+k*Inches(0.70),
       SW-Inches(0.45), Inches(0.62),
       sz=Pt(16), color=DARK_GRAY)

card_body(sl, CL+SW+Inches(0.2), BT5, PW, BH5)
tx(sl, "Key Properties",
   CL+SW+Inches(0.38), BT5+Inches(0.12), PW-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=GREEN)
blist(sl,
    ["4.7M parameters - 8.8× fewer than MLP",
     "Resolution-invariant: weights don't depend on grid size",
     "  → Evaluates at 128×128 zero-shot, no retraining",
     "Spectral inductive bias: Darcy pressure is low-frequency dominated",
     "We also plan to train a larger config:",
     "  → width=128, modes=20, 6 layers - 78.8M params",
     "  → Will study accuracy vs parameter-count trade-off"],
    CL+SW+Inches(0.2), BT5+Inches(0.58), PW, BH5-Inches(0.65),
    sz=Pt(16))

# ══════════════════════════════════════════════════════════════════════════════
# S6 Evaluation Plan
# layout: [primary metric banner] 0.95"
#         [4 metric cards] 2.75"
#         [literature note] 1.5"
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Evaluation Plan", "Metrics, expected results, and planned figures")
footer(sl, 6)

# Primary metric banner  t=1.55 h=0.95 b=2.5
box(sl, CL, CT, CW, Inches(0.95), fill=LIGHT_BLUE, border=ACCENT_BLUE, bw=Pt(2))
tx(sl, "Primary Metric - Relative L2 Error (physical units)",
   CL+Inches(0.25), CT+Inches(0.1), CW-Inches(0.5), Inches(0.4),
   sz=Pt(20), bold=True, color=DARK_BLUE)
tx(sl, "‖ prediction − truth ‖₂ / ‖ truth ‖₂   computed after denormalizing to physical pressure units",
   CL+Inches(0.25), CT+Inches(0.54), CW-Inches(0.5), Inches(0.34),
   sz=Pt(15), color=DARK_GRAY)

# 4 metric cards  t=2.6 h=2.75 b=5.35
MT = CT + Inches(0.95) + Inches(0.1)
MH = Inches(2.75)
MC4W = (CW - Inches(0.6)) / 4
for j, (title, desc, col) in enumerate([
    ("Relative L2 Error",
     "Mean / median / std /\nmin / max across all\n200 test samples",
     ACCENT_BLUE),
    ("Zero-Shot Super-Res",
     "FNO at native 128×128\nvs real PDEBench ground\ntruth - no retraining",
     GREEN),
    ("Inference Speed",
     "Wall-clock time vs FDM\nnumerical solver -\nmeasured, not asserted",
     ORANGE),
    ("Qualitative Plots",
     "κ input | target u |\nprediction | absolute\nerror maps",
     MID_BLUE),
]):
    CX4 = CL + j*(MC4W+Inches(0.2))
    card_hdr(sl, title, CX4, MT, MC4W, fill=col)
    card_body(sl, CX4, MT+Inches(0.48), MC4W, MH-Inches(0.48))
    tx(sl, desc, CX4+Inches(0.18), MT+Inches(0.64), MC4W-Inches(0.36), MH-Inches(0.72),
       sz=Pt(16), color=DARK_GRAY)

# Literature note  t=5.45 h=1.6 b=7.05
LT = MT + MH + Inches(0.1)
card_body(sl, CL, LT, CW, Inches(1.6))
tx(sl, "Literature Reference & Expected Range",
   CL+Inches(0.18), LT+Inches(0.12), CW-Inches(0.36), Inches(0.38),
   sz=Pt(18), bold=True, color=DARK_BLUE)
tx(sl, "Li et al. 2021 (original FNO paper) reports ~0.01–0.02 relative L2 on Darcy flow. "
       "Their dataset uses smooth, continuous log-permeability (κ = exp(GRF)). "
       "PDEBench uses piecewise-constant κ ∈ {0.1, 1.0} - sharper interfaces, a harder variant. "
       "Some deviation from the literature number is expected and will be discussed.",
   CL+Inches(0.18), LT+Inches(0.58), CW-Inches(0.36), Inches(0.95),
   sz=Pt(15), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# S7 Expected Figures
# 5 columns, full height
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Expected Figures & Tables", "Planned deliverables for Stage 2 and Stage 3 reports")
footer(sl, 7)

F5W = (CW - Inches(0.8)) / 5
FH  = CB - CT  # full content height
HDR_H  = Inches(0.68)  # taller header so long titles don't clip
DESC_H = Inches(2.4)   # enough for ~8 wrapped lines at 15pt
PH_T   = CT + HDR_H + Inches(0.1) + DESC_H + Inches(0.1)
PH_H   = CB - PH_T - Inches(0.05)
for j, (num, title, desc, sub, col) in enumerate([
    ("1", "Error Histograms",
     "Per-sample relative L2 for FNO original, FNO improved, and MLP baseline.\nReveals mean, median, and tail behaviour.",
     "X-axis: relative L2 error\nY-axis: sample count\nOne histogram per model, overlaid for direct comparison.",
     ACCENT_BLUE),
    ("2", "Sample Predictions",
     "Side-by-side: κ input, target u, FNO prediction, absolute error.\nShown for several representative test cases.",
     "4-panel grid per sample.\nSelected cases include easy, medium, and hard permeability patterns.",
     MID_BLUE),
    ("3", "Comparison Table",
     "Rel L2 (mean/median), parameter count, and inference time for all three models in one table.",
     "Rows: MLP, FNO original, FNO improved.\nColumns: Rel L2 mean, Rel L2 median, params, ms/sample.",
     GREEN),
    ("4", "Zero-Shot Super-Res",
     "FNO evaluated at native 128x128 vs PDEBench ground truth. No retraining - validates resolution-invariance.",
     "MLP cannot do this (fixed I/O size).\nFNO tested at 64x64 (train res) and 128x128 (zero-shot).",
     ORANGE),
    ("5", "Efficiency Scatter",
     "Rel L2 vs parameter count - one point per model. Shows the accuracy vs model-size trade-off.",
     "We expect FNO lower-left (fewer params, lower error).\nMLP expected upper-right - more parameters, likely higher error.",
     DARK_BLUE),
]):
    CX5 = CL + j*(F5W+Inches(0.2))
    box(sl, CX5, CT, F5W, HDR_H, fill=col)
    tx(sl, f"{num}.  {title}", CX5+Inches(0.12), CT+Inches(0.1),
       F5W-Inches(0.24), HDR_H-Inches(0.12), sz=Pt(13), bold=True, color=WHITE)
    card_body(sl, CX5, CT+HDR_H, F5W, FH-HDR_H)
    box(sl, CX5, CT+HDR_H, F5W, Inches(0.06), fill=col)
    tx(sl, desc, CX5+Inches(0.15), CT+HDR_H+Inches(0.1), F5W-Inches(0.3), DESC_H,
       sz=Pt(15), color=DARK_GRAY)
    tx(sl, sub, CX5+Inches(0.15), PH_T, F5W-Inches(0.3), PH_H,
       sz=Pt(14), color=MED_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# S8 Work Plan
# 3 columns, full height
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Work Plan", "Three-stage project timeline - Stage 1 complete")
footer(sl, 8)

C3W2 = (CW - Inches(0.4)) / 3
for j, (stg, title, date, tasks, col, status) in enumerate([
    ("Stage 1", "Project Proposal", "Aug 27, 2026",
     ["Confirm PDEBench as dataset (source verified)",
      "Finalize MLP and FNO architecture design",
      "Define evaluation methodology and metrics",
      "Submit proposal PDF and presentation"],
     ACCENT_BLUE, "✓  DONE"),
    ("Stage 2", "Implementation & Interim Results", "Sep 8, 2026",
     ["Set up data download and preprocessing pipeline",
      "Implement and train MLP baseline",
      "Implement and train FNO (original config)",
      "Gather preliminary results and plots"],
     MID_BLUE, "UPCOMING"),
    ("Stage 3", "Final Report, Code & Viva", "Nov 13, 2026",
     ["Train improved FNO (width=128, modes=20)",
      "Zero-shot super-resolution at 128×128",
      "Inference speed benchmark vs FDM solver",
      "Final report, code submission, viva"],
     GREEN, "PLANNED"),
]):
    CX3b = CL + j*(C3W2+Inches(0.2))
    # Stage header
    box(sl, CX3b, CT, C3W2, Inches(0.55), fill=col)
    tx(sl, stg, CX3b+Inches(0.15), CT+Inches(0.1), Inches(1.6), Inches(0.35),
       sz=Pt(17), bold=True, color=WHITE)
    tx(sl, status, CX3b+C3W2-Inches(1.5), CT+Inches(0.12), Inches(1.4), Inches(0.3),
       sz=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Sub-header
    box(sl, CX3b, CT+Inches(0.55), C3W2, Inches(0.65), fill=LIGHT_BLUE)
    tx(sl, title, CX3b+Inches(0.15), CT+Inches(0.61), C3W2-Inches(0.3), Inches(0.3),
       sz=Pt(15), bold=True, color=DARK_BLUE)
    tx(sl, f"Due: {date}", CX3b+Inches(0.15), CT+Inches(0.91), C3W2-Inches(0.3), Inches(0.25),
       sz=Pt(13), color=MED_GRAY)
    # Tasks
    card_body(sl, CX3b, CT+Inches(1.2), C3W2, CB-CT-Inches(1.2))
    done = status == "✓  DONE"
    for k, task in enumerate(tasks):
        tx(sl, ("✓  " if done else "→  ") + task,
           CX3b+Inches(0.18), CT+Inches(1.38)+k*Inches(1.0),
           C3W2-Inches(0.36), Inches(0.88),
           sz=Pt(16), color=MID_BLUE if done else DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# S9 Team & AI Declaration
# 4 member cards + AI banner
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=OFF_WHITE)
header(sl, "Individual Contributions", "Team PINNacles - Group of 4")
footer(sl, 9)

M4W = (CW - Inches(0.6)) / 4
MC_H = Inches(3.95)
for j, (name, tasks, col) in enumerate([
    ("Aryamann\nSrivastava",
     ["Model implementation (FNO and MLP)",
      "Training and evaluation pipelines",
      "Data download and preprocessing",
      "Benchmarks and report drafting"],
     ACCENT_BLUE),
    ("Varun\nSathaye",
     ["Literature review",
      "Baseline model tuning",
      "Error analysis",
      "Formatting deliverables"],
     MID_BLUE),
    ("Atishay\nJain",
     ["Data visualization and plotting",
      "Error distribution figures",
      "Field comparison plots",
      "Slide preparation"],
     GREEN),
    ("Vedant S.\nTiwari",
     ["Zero-shot super-res testing",
      "Speed benchmarking",
      "Cross-checking results",
      "Proofreading final report"],
     ORANGE),
]):
    CXM = CL + j*(M4W+Inches(0.2))
    box(sl, CXM, CT, M4W, Inches(0.85), fill=col)
    tx(sl, name, CXM, CT+Inches(0.08), M4W, Inches(0.7),
       sz=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card_body(sl, CXM, CT+Inches(0.85), M4W, MC_H-Inches(0.85))
    for k, t in enumerate(tasks):
        tx(sl, f"•  {t}", CXM+Inches(0.15), CT+Inches(1.02)+k*Inches(0.75),
           M4W-Inches(0.3), Inches(0.68), sz=Pt(15), color=DARK_GRAY)

# AI declaration  t=5.6 h=1.45 b=7.05
AT = CT + MC_H + Inches(0.1)
box(sl, CL, AT, CW, Inches(1.35), fill=LIGHT_BLUE, border=ACCENT_BLUE, bw=Pt(1.5))
tx(sl, "AI Tool Use Declaration",
   CL+Inches(0.25), AT+Inches(0.12), Inches(3.5), Inches(0.38),
   sz=Pt(17), bold=True, color=DARK_BLUE)
tx(sl, "Claude / ChatGPT used for conceptual understanding, coding assistance, debugging, and report drafting. "
       "All output inspected, verified, and fully understood by the team before submission - "
       "in accordance with AE646's academic integrity policy.",
   CL+Inches(0.25), AT+Inches(0.55), CW-Inches(0.5), Inches(0.72),
   sz=Pt(15), color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# S10 Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
box(sl, 0, 0, W, H, fill=DARK_BLUE)
box(sl, 0, Inches(1.12), W, Inches(0.07), fill=ACCENT_BLUE)
tx(sl, "Summary", Inches(0.6), Inches(0.18), W-Inches(1.2), Inches(0.85),
   sz=Pt(36), bold=True, color=WHITE)

rows = [
    ("Problem",    "2D Darcy flow - learn G: κ → u to replace repeated PDE solves"),
    ("Dataset",    "Real PDEBench 2D Darcy (β=1.0)  |  900 / 100 / 200 split at 64×64"),
    ("Baseline",   "MLP (42M params) - no spatial inductive bias, fixed resolution"),
    ("SciML",      "FNO (4.7M params) - spectral convolutions, resolution-invariant"),
    ("Extended",   "Larger FNO (78.8M params) - accuracy vs parameter-count trade-off"),
    ("Evaluation", "Rel L2 error  |  zero-shot super-res at 128×128  |  speed vs FDM"),
]
RH = Inches(0.84)
for j, (label, text) in enumerate(rows):
    RT = Inches(1.32) + j*RH
    box(sl, Inches(0.5), RT+Inches(0.05), Inches(2.0), RH-Inches(0.1), fill=ACCENT_BLUE)
    tx(sl, label, Inches(0.5), RT+Inches(0.12), Inches(2.0), RH-Inches(0.15),
       sz=Pt(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(sl, text, Inches(2.7), RT+Inches(0.14), W-Inches(3.2), RH-Inches(0.18),
       sz=Pt(18), color=LIGHT_BLUE)

box(sl, 0, H-Inches(0.62), W, Inches(0.62), fill=MID_BLUE)
tx(sl, "Thank you  -  Questions welcome",
   Inches(0.5), H-Inches(0.55), W-Inches(1.0), Inches(0.46),
   sz=Pt(19), color=WHITE, align=PP_ALIGN.CENTER)

# ── save ──────────────────────────────────────────────────────────────────────
import os
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "PINNacles_Stage1_Presentation.pptx")
prs.save(OUT)
print(f"Saved → {OUT}")
